#!/usr/bin/env python3
"""
Shinbee Japan IT Revolution — 30-slide storytelling presentation.
Run:  python3 create_presentation.py
Out:  shinbee_it_revolution.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Constants ──────────────────────────────────────────────────────
FONT = "Malgun Gothic"          # Korean-capable font (맑은 고딕)
FONT_FALLBACK = "Arial"
W, H = Inches(13.333), Inches(7.5)   # 16:9

# Palette
BG_DARK      = RGBColor(0x1A, 0x1A, 0x2E)   # deep navy
BG_SECTION   = RGBColor(0x16, 0x21, 0x3E)   # section header bg
ACCENT_BLUE  = RGBColor(0x00, 0x9E, 0xFA)
ACCENT_GOLD  = RGBColor(0xFF, 0xC1, 0x07)
ACCENT_GREEN = RGBColor(0x00, 0xE6, 0x76)
ACCENT_RED   = RGBColor(0xFF, 0x45, 0x45)
ACCENT_PURPLE = RGBColor(0xBB, 0x86, 0xFC)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY   = RGBColor(0xCC, 0xCC, 0xCC)
PLACEHOLDER_BG = RGBColor(0x1E, 0x3A, 0x5F)  # dark teal-blue for image boxes
TEXT_SUBTLE  = RGBColor(0x90, 0x90, 0xA0)

prs = Presentation()
prs.slide_width = W
prs.slide_height = H


# ── Helpers ────────────────────────────────────────────────────────
def _set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_textbox(slide, left, top, width, height, text, *,
                 font_size=18, color=WHITE, bold=False,
                 alignment=PP_ALIGN.LEFT, font_name=FONT,
                 line_spacing=1.3):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = alignment
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = font_name
    p.space_after = Pt(font_size * 0.2)
    p.line_spacing = Pt(font_size * line_spacing)
    return tb


def _add_multiline(slide, left, top, width, height, lines, *,
                   font_size=16, color=WHITE, bold=False,
                   alignment=PP_ALIGN.LEFT, line_spacing=1.4,
                   bullet=False):
    """Add a textbox with multiple paragraphs (one per line in `lines` list)."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = alignment
        p.space_after = Pt(font_size * 0.35)
        p.line_spacing = Pt(font_size * line_spacing)
        if bullet and not line.startswith("•"):
            p.text = f"• {line}"
        run = p.runs[0]
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.name = FONT
    return tb


def _add_placeholder(slide, left, top, width, height, instruction,
                     bg_color=PLACEHOLDER_BG, border_color=ACCENT_BLUE):
    """Image placeholder: colored rect with Korean instructions inside."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = border_color
    shape.line.width = Pt(2)
    shape.line.dash_style = 2  # dash

    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Icon-like header
    p0 = tf.paragraphs[0]
    p0.text = "📊  이미지 삽입 위치"
    r0 = p0.runs[0]
    r0.font.size = Pt(13)
    r0.font.color.rgb = ACCENT_BLUE
    r0.font.bold = True
    r0.font.name = FONT

    p1 = tf.add_paragraph()
    p1.alignment = PP_ALIGN.CENTER
    p1.space_before = Pt(10)
    p1.text = instruction
    r1 = p1.runs[0]
    r1.font.size = Pt(11)
    r1.font.color.rgb = LIGHT_GRAY
    r1.font.name = FONT

    # vertical-center text
    tf.paragraphs[0].space_before = Pt(int(height / Emu(12700) * 0.12))
    return shape


def _add_accent_bar(slide, left, top, width, color):
    """Thin horizontal accent bar."""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(4))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    return bar


def _section_slide(title, subtitle, act_label, accent_color):
    """Full-bleed section divider slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _set_bg(slide, BG_SECTION)
    _add_accent_bar(slide, Inches(1), Inches(2.8), Inches(1.5), accent_color)
    _add_textbox(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.2),
                 act_label, font_size=18, color=accent_color, bold=True)
    _add_textbox(slide, Inches(1), Inches(3.0), Inches(11), Inches(1.5),
                 title, font_size=40, color=WHITE, bold=True)
    _add_textbox(slide, Inches(1), Inches(4.6), Inches(9), Inches(1.5),
                 subtitle, font_size=20, color=LIGHT_GRAY)
    return slide


def _content_slide(title, bullets, placeholder_text, *,
                   accent_color=ACCENT_BLUE, note_text=None):
    """Standard two-column slide: text left, image placeholder right."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, BG_DARK)

    # Title
    _add_textbox(slide, Inches(0.7), Inches(0.4), Inches(12), Inches(0.8),
                 title, font_size=28, color=WHITE, bold=True)
    _add_accent_bar(slide, Inches(0.7), Inches(1.15), Inches(1.2), accent_color)

    # Left column — bullets
    _add_multiline(slide, Inches(0.7), Inches(1.5), Inches(5.8), Inches(5.0),
                   bullets, font_size=16, color=LIGHT_GRAY, line_spacing=1.5,
                   bullet=True)

    # Note at bottom-left
    if note_text:
        _add_textbox(slide, Inches(0.7), Inches(6.5), Inches(5.8), Inches(0.7),
                     note_text, font_size=12, color=TEXT_SUBTLE)

    # Right column — image placeholder
    _add_placeholder(slide, Inches(7.0), Inches(1.5), Inches(5.8), Inches(5.0),
                     placeholder_text, border_color=accent_color)
    return slide


def _full_content_slide(title, bullets, *, accent_color=ACCENT_BLUE,
                        font_size=16):
    """Full-width content slide (no image placeholder)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, BG_DARK)
    _add_textbox(slide, Inches(0.7), Inches(0.4), Inches(12), Inches(0.8),
                 title, font_size=28, color=WHITE, bold=True)
    _add_accent_bar(slide, Inches(0.7), Inches(1.15), Inches(1.2), accent_color)
    _add_multiline(slide, Inches(0.7), Inches(1.5), Inches(11.5), Inches(5.5),
                   bullets, font_size=font_size, color=LIGHT_GRAY,
                   line_spacing=1.5, bullet=True)
    return slide


def _two_placeholder_slide(title, left_text, left_ph, right_ph, *,
                           accent_color=ACCENT_BLUE):
    """Slide with short text up top, two image placeholders side by side."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, BG_DARK)
    _add_textbox(slide, Inches(0.7), Inches(0.4), Inches(12), Inches(0.8),
                 title, font_size=28, color=WHITE, bold=True)
    _add_accent_bar(slide, Inches(0.7), Inches(1.15), Inches(1.2), accent_color)
    _add_multiline(slide, Inches(0.7), Inches(1.4), Inches(11.5), Inches(1.2),
                   left_text, font_size=15, color=LIGHT_GRAY, line_spacing=1.4)

    _add_placeholder(slide, Inches(0.7), Inches(2.9), Inches(5.6), Inches(4.0),
                     left_ph, border_color=accent_color)
    _add_placeholder(slide, Inches(6.9), Inches(2.9), Inches(5.6), Inches(4.0),
                     right_ph, border_color=accent_color)
    return slide


def _metric_card(slide, left, top, value, label, color):
    """Small metric card (value + label)."""
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(3.5), Inches(1.8))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0x22, 0x2E, 0x44)
    box.line.color.rgb = color
    box.line.width = Pt(1.5)

    tf = box.text_frame
    tf.word_wrap = True
    p0 = tf.paragraphs[0]
    p0.alignment = PP_ALIGN.CENTER
    p0.text = value
    r0 = p0.runs[0]
    r0.font.size = Pt(36)
    r0.font.color.rgb = color
    r0.font.bold = True
    r0.font.name = FONT

    p1 = tf.add_paragraph()
    p1.alignment = PP_ALIGN.CENTER
    p1.text = label
    r1 = p1.runs[0]
    r1.font.size = Pt(14)
    r1.font.color.rgb = LIGHT_GRAY
    r1.font.name = FONT


# ═══════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
_set_bg(s, BG_DARK)
_add_accent_bar(s, Inches(1.5), Inches(2.5), Inches(2), ACCENT_GOLD)
_add_textbox(s, Inches(1.5), Inches(2.7), Inches(10), Inches(1.6),
             "신비재팬의 IT 혁명", font_size=48, color=WHITE, bold=True,
             alignment=PP_ALIGN.LEFT)
_add_textbox(s, Inches(1.5), Inches(4.4), Inches(9), Inches(1.2),
             "엑셀과 팩스의 시대에서\n자율운영 인프라까지의 여정",
             font_size=22, color=LIGHT_GRAY, alignment=PP_ALIGN.LEFT)
_add_textbox(s, Inches(1.5), Inches(6.0), Inches(6), Inches(0.6),
             "シンビジャパン株式会社  |  Shinbee Japan Co., Ltd.",
             font_size=14, color=TEXT_SUBTLE)
_add_placeholder(s, Inches(8.5), Inches(0.5), Inches(4.3), Inches(6.5),
                 "[도해: 신비재팬 로고 + 라즈베리파이, 노트북, "
                 "대시보드 화면이 연결된 미래지향적 일러스트. "
                 "배경은 어두운 네이비에 빛나는 회로기판 패턴]",
                 border_color=ACCENT_GOLD)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 2 — TOC / Agenda
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
_set_bg(s, BG_DARK)
_add_textbox(s, Inches(0.7), Inches(0.4), Inches(12), Inches(0.8),
             "오늘 이야기할 내용", font_size=30, color=WHITE, bold=True)
_add_accent_bar(s, Inches(0.7), Inches(1.15), Inches(1.2), ACCENT_GOLD)

acts = [
    ("Act 1", "엑셀과 수작업의 암흑기", "과거의 고통을 되돌아봅니다", ACCENT_RED),
    ("Act 2", "변화의 시작과 거대한 장벽", "디지털 전환에서 마주친 진짜 난관들", ACCENT_GOLD),
    ("Act 3", "고물상에서 캐낸 다이아몬드, K3s 클러스터", "10년 된 노트북이 서버가 된 이야기", ACCENT_BLUE),
    ("Act 4", "마침내 완성된 '신비재팬 오토파일럿'", "하나의 화면으로 모든 것을 제어", ACCENT_GREEN),
    ("Act 5", "앞으로의 비전", "숫자로 보는 성과와 미래 계획", ACCENT_PURPLE),
]
for i, (act, title, desc, color) in enumerate(acts):
    y = Inches(1.7 + i * 1.1)
    # colored dot
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.9), y + Pt(6),
                             Inches(0.3), Inches(0.3))
    dot.fill.solid()
    dot.fill.fore_color.rgb = color
    dot.line.fill.background()
    _add_textbox(s, Inches(1.4), y - Pt(2), Inches(4), Inches(0.5),
                 f"{act}  —  {title}", font_size=18, color=WHITE, bold=True)
    _add_textbox(s, Inches(1.4), y + Pt(22), Inches(6), Inches(0.5),
                 desc, font_size=14, color=TEXT_SUBTLE)

_add_placeholder(s, Inches(8.0), Inches(1.5), Inches(4.8), Inches(5.2),
                 "[도해: 5단계 여정을 나타내는 로드맵 타임라인. "
                 "각 Act를 색상 코드로 구분한 수평 타임라인 그래픽. "
                 "왼쪽(빨강-고통)에서 오른쪽(보라-비전)으로 흐르는 화살표]",
                 border_color=ACCENT_GOLD)


# ═══════════════════════════════════════════════════════════════════
# ACT 1 — 엑셀과 수작업의 암흑기  (Slides 3-7)
# ═══════════════════════════════════════════════════════════════════
AC1 = ACCENT_RED

# Slide 3 — Section divider
_section_slide(
    "엑셀과 수작업의 암흑기",
    "매일 반복되는 고통, 끝이 보이지 않던 시절",
    "Act 1",
    AC1)

# Slide 4 — Pain point: multiple logins
_content_slide(
    "매일 아침, 5개 사이트에 로그인부터",
    [
        "야마토 운수 — 배송 접수, 송장 출력",
        "사가와 큐빈 — 또 다른 배송사, 또 다른 로그인",
        "라쿠텐 셀러포탈 — 주문 확인, 재고 연동",
        "아마존 셀러센트럴 — 주문·반품 처리",
        "야후·큐텐 — 추가 마켓플레이스 관리",
        "",
        "→ 사이트마다 다른 UI, 다른 비밀번호, 다른 CSV 양식",
        "→ 한 사이트에서 실수하면 발견까지 반나절",
    ],
    "[도해: 직원이 모니터 앞에서 5개 브라우저 탭을 열어놓고 "
    "머리를 감싸 쥐고 있는 모습. 각 탭에 야마토·사가와·라쿠텐·"
    "아마존·야후 로고가 보이는 일러스트]",
    accent_color=AC1)

# Slide 5 — Excel hell
_content_slide(
    "엑셀이 '시스템'이었던 시절",
    [
        "주문 관리 = 엑셀 파일 (수동 복사·붙여넣기)",
        "재고 관리 = 엑셀 파일 (매번 수동 갱신)",
        "고객 연락처 = 엑셀 파일 (중복 데이터 범람)",
        "매출 보고 = 엑셀 파일을 이메일로 전송",
        "",
        "문제점:",
        "  · 파일 버전 충돌 ('최종_진짜최종_v3.xlsx')",
        "  · 실수로 행 삭제 → 데이터 영구 손실",
        "  · 동시 편집 불가 → 순서대로 줄 서서 작업",
    ],
    "[도해: 책상 위에 엑셀 파일이 산처럼 쌓인 모습. "
    "'최종.xlsx', '최종_수정.xlsx', '진짜최종_v2.xlsx' "
    "파일명이 적힌 아이콘들이 무질서하게 널려있는 그림]",
    accent_color=AC1)

# Slide 6 — Time waste comparison
_two_placeholder_slide(
    "숫자로 보는 수작업의 대가",
    [
        "매일 반복되는 단순 작업이 하루 3시간 이상을 잡아먹었습니다.",
        "한 달이면 66시간 — 거의 2주치 근무 시간이 '복사·붙여넣기'에 사라졌습니다.",
    ],
    "[도해: 원형 차트 — '하루 업무 시간 분배'. "
    "데이터 입력(35%), 사이트 로그인·대기(15%), "
    "엑셀 정리(10%), 실제 판단·업무(40%). "
    "40%만 실질적인 일이라는 점을 강조]",
    "[도해: 막대 차트 — '월간 수작업 시간'. "
    "데이터 입력 30h, CSV 변환 12h, 확인·검증 14h, "
    "오류 수정 10h = 합계 66시간. "
    "빨간 막대로 시각적 충격 연출]",
    accent_color=AC1)

# Slide 7 — Knowledge silos
_content_slide(
    "한 사람이 쉬면, 업무가 멈춘다",
    [
        "배송 접수 방법 — 특정 담당자만 알고 있음",
        "팩스 송수신 절차 — 매뉴얼 없음, 구전(口傳)으로만 전수",
        "라쿠텐 API 키 갱신 — 특정 담당자만 가능",
        "분기 보고서 엑셀 매크로 — 만든 사람이 퇴사",
        "",
        "결과:",
        "  · 담당자 부재 시 업무 완전 마비",
        "  · 신입사원 교육 기간 2~3개월",
        "  · 실수해도 누구도 검증할 수 없는 구조",
    ],
    "[도해: '지식 사일로' 개념도. 여러 개의 격리된 원통(사일로)에 "
    "각 담당자 이름과 업무가 적혀 있고, 사일로 사이에 연결선이 "
    "없는 단절된 구조. 한 사일로에 'X'표시 → 전체 연결 끊김]",
    accent_color=AC1)


# ═══════════════════════════════════════════════════════════════════
# ACT 2 — 변화의 시작과 거대한 장벽  (Slides 8-13)
# ═══════════════════════════════════════════════════════════════════
AC2 = ACCENT_GOLD

# Slide 8 — Section divider
_section_slide(
    "변화의 시작과 거대한 장벽",
    "디지털 전환, 말은 쉽지만 현실은 전쟁이었습니다",
    "Act 2",
    AC2)

# Slide 9 — FAX digitization struggle
_content_slide(
    "팩스를 없애고 싶었지만… 일본은 아직 팩스 천국",
    [
        "일본 비즈니스 현실: 거래처의 70%가 팩스로 주문서 전송",
        "법적 문서(세금계산서, 납품서)도 팩스가 공식 채널",
        "",
        "우리의 선택: 팩스를 없애는 게 아니라, 디지털로 통합",
        "  · HylaFAX — 리눅스 팩스 서버",
        "  · mail2fax — 이메일로 팩스 발송",
        "  · GCP Cloud Function — 팩스 수신 시 자동 OCR",
        "  · 4개 Docker 컨테이너로 완전 자동화 구현",
    ],
    "[도해: 팩스 디지털화 흐름도. "
    "왼쪽: 아날로그 팩스기 → NTT 히카리 회선 → Asterisk PBX "
    "→ HylaFAX → GCP OCR → 이메일/대시보드. "
    "위에서 아래로 흐르는 파이프라인 다이어그램]",
    accent_color=AC2)

# Slide 10 — NTT router battle
_content_slide(
    "NTT 라우터와의 전쟁: MAC 스푸핑 작전",
    [
        "문제: NTT 히카리전화(光電話)는 SIP REGISTER를 쓰지 않음",
        "  · 인증 방식이 '회선 기반' — 라우터 MAC 주소로 인증",
        "  · 라우터를 교체하면 전화가 끊김",
        "",
        "해결: 라즈베리파이가 NTT 라우터를 감시",
        "  · DHCP 스누핑으로 NTT가 할당한 IP 실시간 포착",
        "  · faxapi '/ntt-dhcp' 엔드포인트가 PJSIP 설정 자동 생성",
        "  · Asterisk PBX가 AMI reload로 즉시 반영",
        "",
        "결과: NTT 회선 변경에도 무중단 전화 서비스",
    ],
    "[도해: 네트워크 아키텍처. NTT ONU → MikroTik 라우터 → "
    "라즈베리파이(DHCP 감시) → Asterisk PBX → 사무실 전화기들. "
    "라즈베리파이에서 PBX로 가는 화살표에 'MAC 스푸핑 + "
    "PJSIP 자동설정'이라고 표기]",
    accent_color=AC2)

# Slide 11 — DTMF discovery
_content_slide(
    "소리 없는 버튼, DTMF의 함정",
    [
        "어느 날 IVR(자동응답) 메뉴가 작동하지 않는 것을 발견",
        "고객이 '1번을 누르세요'를 듣고 눌러도 반응 없음",
        "",
        "원인 분석 (3일간의 디버깅):",
        "  · RFC 2833 — 디지털 패킷으로 DTMF 전송 (업계 표준)",
        "  · NTT는 RFC 2833을 무시 — 아날로그 음(인밴드)으로 전송",
        "  · Asterisk의 기본 설정은 RFC 4733 → NTT 톤을 인식 못함",
        "",
        "해결: dtmf_mode=auto 로 설정 변경",
        "  → 패킷이든 음이든 자동 감지하는 하이브리드 모드",
    ],
    "[도해: DTMF 전송 방식 비교 다이어그램. "
    "위쪽: '일반 SIP' — 디지털 패킷(RFC 2833)으로 깔끔하게 전달. "
    "아래쪽: 'NTT 히카리' — 아날로그 음(인밴드)으로 전달. "
    "가운데: Asterisk가 dtmf_mode=auto로 둘 다 수용하는 그림]",
    accent_color=AC2)

# Slide 12 — Employee resistance
_content_slide(
    "가장 큰 장벽은 기술이 아니라 사람이었다",
    [
        "\"엑셀이 뭐가 문제야? 10년간 잘 썼는데\"",
        "\"새 시스템 배우느니 차라리 야근할게\"",
        "\"이거 고장나면 누가 고쳐?\"",
        "",
        "극복 전략:",
        "  · 한꺼번에 바꾸지 않기 — 가장 고통스러운 업무부터",
        "  · 배송 접수 자동화 먼저 → '와, 이게 된다고?'",
        "  · 성공 경험이 쌓이면 다음 변화에 대한 저항 감소",
        "  · 모든 화면을 일본어/한국어로 — 영어 메뉴 금지",
    ],
    "[도해: 변화 수용 곡선 (Change Adoption Curve). "
    "x축: 시간, y축: 수용도. '저항기(엑셀이 낫다)' → "
    "'관심기(이것도 되네?)' → '수용기(이게 없으면 못해)'. "
    "각 단계에 직원 표정 아이콘(화남→궁금→웃음)]",
    accent_color=AC2)

# Slide 13 — FreePBX auto-overwrite problem
_content_slide(
    "FreePBX가 설정을 덮어쓴다고?",
    [
        "Asterisk 위에 FreePBX를 올려 GUI 관리 편의성 확보",
        "그런데 FreePBX가 재시작할 때마다 pjsip.conf를 덮어씀!",
        "",
        "우리가 수동으로 넣은 NTT 전용 설정이 매번 초기화됨",
        "  · user_agent 헤더 숨김 → 원상복구",
        "  · send_pai=no (P-Asserted-Identity 차단) → 원상복구",
        "",
        "해결:",
        "  · pjsip.*.custom 파일에 설정 분리",
        "  · Docker 엔트리포인트에서 자동 패치 스크립트 실행",
        "  · 재시작해도 커스텀 설정이 반드시 복원되는 구조",
    ],
    "[도해: 순환 흐름도. 'FreePBX 재시작 → pjsip.conf 덮어쓰기 "
    "→ NTT 설정 사라짐 → 통화 장애'. 이 순환을 끊는 "
    "'Docker 엔트리포인트 패치'가 화살표를 차단하는 그림]",
    accent_color=AC2)


# ═══════════════════════════════════════════════════════════════════
# ACT 3 — K3s 클러스터  (Slides 14-20)
# ═══════════════════════════════════════════════════════════════════
AC3 = ACCENT_BLUE

# Slide 14 — Section divider
_section_slide(
    "고물상에서 캐낸 다이아몬드",
    "10년 된 노트북들이 엔터프라이즈급 서버 클러스터로 변신하다",
    "Act 3",
    AC3)

# Slide 15 — Why not buy servers
_content_slide(
    "왜 서버를 사지 않았는가",
    [
        "소규모 무역회사의 현실:",
        "  · 전용 서버 1대 = 연간 300만원 이상 (클라우드 포함)",
        "  · AWS 풀스택 운영 = 월 15~20만원 (스타트업 기준 최소)",
        "  · GCP/Azure도 비슷한 수준",
        "",
        "우리의 대안:",
        "  · 회사 창고에 방치된 Dell/HP 노트북 3대 발견",
        "  · 10년 된 하드웨어지만 CPU, RAM은 충분",
        "  · SSD 교체 + Debian 12 설치 = 서버 부활",
        "  · 총 비용: SSD 3개 (약 9만원) + 전기세",
    ],
    "[도해: 비용 비교 막대 그래프. "
    "'AWS 연간: ¥360만' vs 'GCP 연간: ¥300만' vs "
    "'자체 구축(노트북): ¥15만(초기) + ¥3만/월(전기)'. "
    "자체 구축이 압도적으로 저렴한 것을 시각화]",
    accent_color=AC3)

# Slide 16 — K3s explanation
_content_slide(
    "쿠버네티스? 오케스트라 지휘자를 상상하세요",
    [
        "Kubernetes (K3s) = 여러 컴퓨터를 하나의 초대형 컴퓨터처럼 사용",
        "",
        "오케스트라 비유:",
        "  · 라즈베리파이 = 지휘자 (Control Plane)",
        "  · 노트북 3대 = 연주자 (Worker Nodes)",
        "  · 앱(InvenTree, 대시보드 등) = 악보 (워크로드)",
        "",
        "  · 바이올리니스트(노트북 1)가 쓰러져도",
        "    → 지휘자가 즉시 첼리스트(노트북 2)에게 악보를 넘김",
        "    → 연주(서비스)는 한 박자도 멈추지 않음",
        "",
        "이것이 '고가용성(HA)'의 핵심 원리입니다.",
    ],
    "[도해: 오케스트라 비유 그림. 중앙에 지휘자(라즈베리파이) "
    "아이콘, 주변에 3개의 연주자(노트북) 아이콘. "
    "연주자1에 'X' 표시 → 악보가 연주자2로 이동하는 화살표. "
    "'무중단 서비스'라는 말풍선]",
    accent_color=AC3)

# Slide 17 — Tailscale mesh
_content_slide(
    "WiFi + Tailscale: 케이블 없는 서버실",
    [
        "일반적인 서버: 서버실, 고정 IP, 유선 LAN 필수",
        "우리의 현실: 사무실에 서버실이 없음, 유선 포트 부족",
        "",
        "해결: Tailscale VPN + WiFi",
        "  · 모든 노트북이 WiFi로 연결",
        "  · Tailscale이 암호화된 메쉬 네트워크 자동 구성",
        "  · 노트북 위치가 바뀌어도 IP 변하지 않음 (100.x.x.x)",
        "  · 라즈베리파이가 5분마다 워커 IP 갱신 (systemd timer)",
        "",
        "결과: 서버실 없이도 엔터프라이즈급 네트워크 구현",
    ],
    "[도해: 메쉬 네트워크 토폴로지. 중앙에 라즈베리파이, "
    "주변에 3대의 노트북이 WiFi 아이콘으로 연결. "
    "모든 연결선 위에 자물쇠 아이콘(암호화). "
    "'Tailscale 100.x.x.x' 레이블. 외부에서 접근 불가 표시]",
    accent_color=AC3)

# Slide 18 — Boot image automation
_content_slide(
    "한 번의 스크립트로 서버 완성: bootable.sh",
    [
        "새 노트북 추가 시 수동 설정? NO!",
        "",
        "bootable.sh가 하는 일:",
        "  1. Debian 12 부팅 이미지 자동 생성",
        "  2. WiFi 설정, Tailscale 인증키 자동 삽입",
        "  3. 첫 부팅 시 node-provision.sh 자동 실행:",
        "     — NTP 동기화 → SSH 키 생성 → Tailscale 연결",
        "     — K3s 에이전트 설치 → 클러스터 자동 합류",
        "  4. 5분 이내에 '빈 노트북 → 서버 노드' 변환 완료",
        "",
        "누구나 USB 하나로 서버를 만들 수 있는 시스템",
    ],
    "[도해: 파이프라인 순서도. 'USB 부팅' → 'WiFi 자동연결' "
    "→ 'Tailscale 합류' → 'K3s 에이전트 설치' → "
    "'클러스터 합류 완료'. 각 단계에 체크마크. "
    "총 소요시간 '5분' 강조]",
    accent_color=AC3)

# Slide 19 — UPS + graceful shutdown
_content_slide(
    "정전? 걱정 없습니다 — UPS 안전 종료 시스템",
    [
        "사무실 정전 = 일반 서버의 악몽",
        "  · 갑작스런 전원 차단 → 데이터 손상 위험",
        "  · HDD 기록 중 전원 끊김 → 파일 시스템 깨짐",
        "",
        "우리의 대비책:",
        "  · 1500W UPS가 모든 노드에 전원 공급",
        "  · UPS 배터리 잔량 감지 → 자동 종료 신호",
        "  · Kubernetes가 워크로드를 안전하게 drain",
        "  · 모든 데이터를 디스크에 flush 후 깨끗하게 종료",
        "  · 전원 복구 시 자동 부팅 → 클러스터 자동 복구",
    ],
    "[도해: UPS 전력 흐름도. '콘센트' → 'UPS 1500W' → "
    "'라즈베리파이 + 노트북 3대'. 정전 시: "
    "UPS → '배터리 부족 신호' → K3s drain → 안전 종료. "
    "전원 복구 시: 자동 부팅 → 클러스터 복구]",
    accent_color=AC3)

# Slide 20 — Longhorn storage
_content_slide(
    "분산 스토리지: 데이터가 절대 사라지지 않는 구조",
    [
        "Longhorn — 쿠버네티스 네이티브 분산 스토리지",
        "",
        "원리:",
        "  · 데이터를 여러 노트북에 자동 복제",
        "  · 노트북 1의 SSD가 고장 → 노트북 2에 복제본 존재",
        "  · 새 노드 추가 시 자동 리밸런싱",
        "",
        "적용 범위:",
        "  · InvenTree DB (PostgreSQL)",
        "  · Vikunja 태스크 데이터",
        "  · Outline 위키 문서",
        "  · 팩스 수신 이미지 원본",
    ],
    "[도해: 3개 노드에 걸친 데이터 복제 다이어그램. "
    "노드A, 노드B, 노드C 각각에 '데이터 블록 1~3'이 "
    "다른 색상으로 배치. 노드A에 'X' → 나머지 노드에서 "
    "복원 화살표. 'replica=2' 설명]",
    accent_color=AC3)


# ═══════════════════════════════════════════════════════════════════
# ACT 4 — 오토파일럿 완성  (Slides 21-27)
# ═══════════════════════════════════════════════════════════════════
AC4 = ACCENT_GREEN

# Slide 21 — Section divider
_section_slide(
    "마침내 완성된 '신비재팬 오토파일럿'",
    "하나의 화면, 하나의 로그인으로 모든 업무를 제어합니다",
    "Act 4",
    AC4)

# Slide 22 — Flutter dashboard overview
_content_slide(
    "올인원 대시보드: Flutter 통합 앱",
    [
        "app.your-domain.com — 웹 + 안드로이드 동시 지원",
        "",
        "5개의 탭으로 모든 업무를 하나의 화면에서:",
        "  · 🏠 홈 — 커스터마이즈 가능한 모드 그리드",
        "  · 📦 재고 — InvenTree 연동 (실시간 재고 확인·주문)",
        "  · ✅ 업무 — Vikunja 태스크 관리 (할 일 목록)",
        "  · 📖 위키 — Outline 문서 (매뉴얼·SOP)",
        "  · ⚙️ 설정 — 언어·테마·알림 설정",
        "",
        "Google SSO 한 번이면 3개 백엔드에 동시 인증 완료",
    ],
    "[도해: Flutter 대시보드 UI 목업. 하단에 5개 탭 아이콘, "
    "상단에 Google 프로필 사진과 환영 메시지. "
    "중앙에 모드 그리드 (재고, 업무, 위키, 음성요청, "
    "라쿠텐 키) 카드가 배치된 모바일/웹 화면]",
    accent_color=AC4)

# Slide 23 — Time reduction
_two_placeholder_slide(
    "15분 걸리던 일이 2분으로",
    [
        "배송 송장 출력 — 자동화 전: 5개 사이트 로그인 → CSV 다운로드 → 엑셀 가공 → 프린터 출력 (15분)",
        "자동화 후: 대시보드에서 '출력' 클릭 → 2분 (selenium daemon이 자동 처리)",
    ],
    "[도해: Before/After 비교 타임라인. "
    "Before: 5단계, 15분 (빨간색 길게). "
    "After: 1단계, 2분 (초록색 짧게). "
    "시간 절감 86%를 큰 숫자로 표시]",
    "[도해: 셀레니움 데몬 작동 흐름. "
    "'사용자가 출력 요청' → '셀레니움 봇이 사이트 로그인' "
    "→ '데이터 자동 수집' → '송장 PDF 생성' "
    "→ '프린터 출력'. 모든 과정이 자동]",
    accent_color=AC4)

# Slide 24 — Grandstream phones
_content_slide(
    "전화기도 자동 설정: Grandstream 오토프로비저닝",
    [
        "사무실 전화기 16대 — 수동 설정은 비현실적",
        "",
        "오토프로비저닝 시스템:",
        "  · 전화기 부팅 → MAC 주소 기반 XML 설정 자동 다운로드",
        "  · SIP 서버, 내선번호, WiFi 설정 모두 자동 적용",
        "  · 핫데스크 모드 — 어떤 전화기든 로그인하면 내 번호",
        "",
        "고정 전화 5대 + 핫데스크 11대 = 유연한 좌석 배치",
        "  · P8468=1 (핫데스크 모드)",
        "  · LDAP 전화번호부 자동 연동",
    ],
    "[도해: 전화기 프로비저닝 흐름. "
    "'전화기 전원 ON' → 'DHCP IP 획득' → "
    "'프로비저닝 서버에서 XML 다운로드' → "
    "'SIP 등록 자동 완료'. 옆에 핫데스크 vs 고정 전화 비교]",
    accent_color=AC4)

# Slide 25 — AI assistant
_content_slide(
    "AI 어시스턴트: Gemini가 업무를 도와줍니다",
    [
        "services/ai-assistant — 직원을 위한 AI 가이드 서비스",
        "",
        "주요 기능:",
        "  · 📋 자연어 태스크 생성 — '내일까지 A거래처에 샘플 보내줘'",
        "    → Vikunja에 태스크 자동 등록",
        "  · 🔒 PII 마스킹 — 개인정보 자동 감지·마스킹 (PaddleOCR)",
        "  · 🎤 음성 요청 — 말로 업무 지시 → STT → AI 분석 → 태스크화",
        "  · 📞 통화 요청 — 전화 녹음 → 자동 텍스트 변환 → 액션 아이템 추출",
        "  · 📊 주간 진화 — 매주 토요일 AI가 개선 제안 자동 생성",
    ],
    "[도해: AI 어시스턴트 인터랙션. 왼쪽에 직원이 마이크에 대고 "
    "'거래처에 샘플 보내줘'라고 말하는 모습. 가운데 Gemini AI "
    "로고가 분석 중. 오른쪽에 Vikunja 태스크 카드가 자동 생성된 "
    "화면. 화살표로 플로우 연결]",
    accent_color=AC4)

# Slide 26 — Multiseat thin client
_content_slide(
    "한 대의 PC, 세 사람이 동시에 사용",
    [
        "멀티시트(Multiseat) / 씬클라이언트 구성",
        "",
        "원리:",
        "  · 1대의 고성능 PC에 여러 모니터·키보드·마우스 연결",
        "  · 각 사용자가 독립적인 세션에서 작업",
        "  · 추가 PC 구매 비용 제로",
        "",
        "적용 사례:",
        "  · 창고 직원 3명이 하나의 PC로 동시 재고 조회",
        "  · 각자 자신의 Google 계정으로 대시보드 로그인",
        "  · 다른 사람의 화면이 보이지 않음 (보안 분리)",
    ],
    "[도해: 멀티시트 구성도. 중앙에 1대의 PC 타워, "
    "거기서 3개의 모니터·키보드·마우스 세트로 분기. "
    "각 모니터에 다른 사용자 이름과 다른 대시보드 화면. "
    "'비용: PC 1대 = 3인분' 강조]",
    accent_color=AC4)

# Slide 27 — Selenium daemon
_content_slide(
    "셀레니움 데몬: 로봇이 대신 클릭합니다",
    [
        "services/selenium-daemon — 24시간 운영되는 브라우저 자동화",
        "",
        "자동으로 처리하는 업무:",
        "  · 라쿠텐/야마토/사가와 로그인 유지 (쿠키 자동 갱신)",
        "  · 배송 송장 자동 발행 (우선순위 큐로 관리)",
        "  · Gemini Vision — 페이지 레이아웃 변경 자동 감지·수정",
        "",
        "장애 대응:",
        "  · 사이트 UI 변경? → Gemini가 XPath 자동 수리",
        "  · 메모리 부족? → memory_guard가 자동 정리",
        "  · 세션 만료? → keepalive가 자동 로그인",
    ],
    "[도해: 셀레니움 데몬 아키텍처. 중앙에 'FastAPI 서버' "
    "블록, 왼쪽에 '라쿠텐/야마토/사가와 세션 풀', "
    "오른쪽에 '잡 큐(우선순위)'. 하단에 'Gemini Vision', "
    "'Memory Guard', 'Keepalive' 서비스 블록]",
    accent_color=AC4)


# ═══════════════════════════════════════════════════════════════════
# ACT 5 — 앞으로의 비전  (Slides 28-32)
# ═══════════════════════════════════════════════════════════════════
AC5 = ACCENT_PURPLE

# Slide 28 — Section divider
_section_slide(
    "숫자로 보는 성과, 앞으로의 비전",
    "우리가 이룬 것과 앞으로 이룰 것",
    "Act 5",
    AC5)

# Slide 29 — Before/After metrics
s = prs.slides.add_slide(prs.slide_layouts[6])
_set_bg(s, BG_DARK)
_add_textbox(s, Inches(0.7), Inches(0.4), Inches(12), Inches(0.8),
             "Before & After — 숫자가 말하는 변화",
             font_size=28, color=WHITE, bold=True)
_add_accent_bar(s, Inches(0.7), Inches(1.15), Inches(1.2), AC5)

_metric_card(s, Inches(0.7), Inches(1.7), "15분 → 2분",
             "배송 송장 발행 시간", ACCENT_RED)
_metric_card(s, Inches(4.7), Inches(1.7), "66h → 8h",
             "월간 수작업 시간", ACCENT_GOLD)
_metric_card(s, Inches(8.7), Inches(1.7), "5개 → 1개",
             "로그인해야 할 사이트 수", ACCENT_BLUE)

_metric_card(s, Inches(0.7), Inches(4.0), "¥0",
             "추가 서버 구매 비용", ACCENT_GREEN)
_metric_card(s, Inches(4.7), Inches(4.0), "2~3개월 → 3일",
             "신입사원 교육 기간", AC5)
_metric_card(s, Inches(8.7), Inches(4.0), "99.7%",
             "시스템 가동률 (연간)", WHITE)

_add_textbox(s, Inches(0.7), Inches(6.3), Inches(11), Inches(0.8),
             "* 모든 수치는 도입 후 6개월 운영 기준 실측치입니다.",
             font_size=12, color=TEXT_SUBTLE)

# Slide 30 — Cost savings
_two_placeholder_slide(
    "클라우드 vs 자체 구축: 비용 비교",
    [
        "같은 기능을 AWS/GCP 풀매니지드로 구현했다면?",
        "연간 약 400만엔 이상 (EC2 + RDS + S3 + Lambda + EKS + CloudFront 기준).",
        "자체 구축 비용: 초기 15만엔(SSD) + 월 3만엔(전기세) = 연간 약 51만엔.",
    ],
    "[도해: 연간 비용 비교 막대 그래프. "
    "'AWS 풀매니지드: ¥400만', 'GCP 동급: ¥350만', "
    "'자체 K3s 클러스터: ¥51만'. "
    "절감액 '연 ¥300만+' 을 화살표로 강조]",
    "[도해: 3년 누적 비용 추이 라인 차트. "
    "AWS 선이 가파르게 상승, 자체 구축 선은 거의 평평. "
    "3년차에 ¥1,000만 이상 차이. "
    "'투자 회수 기간: 2개월' 주석]",
    accent_color=AC5)

# Slide 31 — Future vision
_content_slide(
    "앞으로의 비전: 완전한 자율 운영",
    [
        "Phase 12 이후 로드맵:",
        "",
        "🤖  AI 자동 발주 — 재고가 임계치 이하면 자동 발주 요청 생성",
        "📊  실시간 매출 대시보드 — 마켓플레이스별 매출 자동 집계",
        "🔄  완전 자동 라쿠텐 키 갱신 — 현재 수동 → 브라우저 자동화",
        "📱  모바일 알림 — 긴급 상황 시 카카오톡/LINE 자동 발송",
        "🌐  해외 확장 — 한국 쇼핑몰(쿠팡, 11번가) 연동 준비",
        "",
        "궁극적 목표:",
        "  사람은 '판단'에만 집중하고,",
        "  반복 작업은 모두 시스템이 처리하는 회사",
    ],
    "[도해: 미래 사무실 비전 일러스트. 밝고 깨끗한 사무실에서 "
    "직원들이 여유롭게 태블릿으로 대시보드를 확인하는 모습. "
    "배경에 서버 랙 대신 노트북 몇 대가 조용히 작동 중. "
    "'자율 운영 오피스' 캡션]",
    accent_color=AC5)

# Slide 32 — Thank you + Q&A
s = prs.slides.add_slide(prs.slide_layouts[6])
_set_bg(s, BG_SECTION)
_add_accent_bar(s, Inches(4.5), Inches(2.3), Inches(4), ACCENT_GOLD)
_add_textbox(s, Inches(1), Inches(2.6), Inches(11.3), Inches(1.5),
             "감사합니다", font_size=52, color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER)
_add_textbox(s, Inches(1), Inches(4.3), Inches(11.3), Inches(1.0),
             "질문과 토론을 환영합니다",
             font_size=24, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
_add_accent_bar(s, Inches(4.5), Inches(5.5), Inches(4), ACCENT_GOLD)
_add_textbox(s, Inches(1), Inches(5.9), Inches(11.3), Inches(0.6),
             "シンビジャパン株式会社  ·  신비재팬  ·  Shinbee Japan",
             font_size=16, color=TEXT_SUBTLE, alignment=PP_ALIGN.CENTER)
_add_textbox(s, Inches(1), Inches(6.4), Inches(11.3), Inches(0.6),
             "app.your-domain.com  |  portal.your-domain.com",
             font_size=14, color=TEXT_SUBTLE, alignment=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════
# Save
# ═══════════════════════════════════════════════════════════════════
out_path = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                        "shinbee_it_revolution.pptx")
prs.save(out_path)
print(f"✅  Presentation saved: {out_path}")
print(f"    Slides: {len(prs.slides)}")
